import os
import logging
import configparser
from whoosh.index import open_dir
from whoosh.qparser import QueryParser
import atexit
from flask import (
    Flask,
    request,
    render_template,
    render_template_string,
    redirect,
    url_for,
    jsonify,
)
import markdown
from markdown.extensions.toc import TocExtension
from markdown.extensions.tables import TableExtension
from markdown.extensions.fenced_code import FencedCodeExtension
from markdown.extensions.codehilite import CodeHiliteExtension
from markdown.extensions.nl2br import Nl2BrExtension
from markdown.extensions.sane_lists import SaneListExtension
from markdown.extensions.smarty import SmartyExtension
from markdown.extensions.wikilinks import WikiLinkExtension
from markdown.extensions.admonition import AdmonitionExtension
from markdown.extensions.attr_list import AttrListExtension
from markdown.extensions.def_list import DefListExtension
from markdown.extensions.footnotes import FootnoteExtension
from markdown.extensions.meta import MetaExtension
from markdown.extensions.abbr import AbbrExtension
from markdown.extensions.legacy_em import LegacyEmExtension
from markdown.extensions.md_in_html import MarkdownInHtmlExtension
import docx
import pptx
from bs4 import BeautifulSoup
import re

# 脚本所在目录
script_dir = os.path.dirname(os.path.abspath(__file__))
# 配置文件路径
config_file = os.path.join(script_dir, "data/config/config.ini")
# 日志文件路径
log_file = os.path.join(script_dir, "data/logs", "searcher.log")
# Whoosh 索引存储的目录
index_dir = os.path.join(script_dir, "data/index_dir")

# 配置解析器
config = configparser.ConfigParser()
config.read(config_file)

# 获取 folders 配置项的值，并将其分割成列表
folder_names = config.get("Folders", "folders").split(",")
# 构建 BASE_DIR 列表，包含 data 文件夹下所有子目录的完整路径
DEL_BASE_DIR = tuple(
    [os.path.join(script_dir, "data", folder_name) for folder_name in folder_names]
)

# 读取 [Logging] 部分的 log_level 配置项
log_level_str = config.get('Logging', 'log_level', fallback='INFO')
log_level = getattr(logging, log_level_str.upper(), logging.INFO)


# 检查日志文件是否存在，如果不存在则创建
log_dir = os.path.dirname(log_file)
if not os.path.exists(log_dir):
    os.makedirs(log_dir)

# 配置日志记录器
logging.basicConfig(
    filename=log_file,
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)


# 定义了一个搜索函数，"dir:xxx"，不使用 jieba 分词的索引结果进行搜索
# 当使用 "dir:xxx" 搜索时，默认给 "xxx" 加上BASE_DIR，以列出该目录下的全部文件，
BASE_DIR = os.path.join(script_dir, "data")

# 配置 Markdown 扩展
markdown_extensions = [
    TocExtension(),  # 支持目录
    TableExtension(),  # 支持表格
    FencedCodeExtension(),  # 支持代码块
    CodeHiliteExtension(),  # 代码高亮
    Nl2BrExtension(),  # 将一个回车视为换行
    SaneListExtension(),  # 支持列表
    SmartyExtension(),  # 支持智能符号
    WikiLinkExtension(),  # 支持 Wiki 链接
    AdmonitionExtension(),  # 支持提示框
    AttrListExtension(),  # 支持属性列表
    DefListExtension(),  # 支持定义列表
    FootnoteExtension(),  # 支持脚注
    MetaExtension(),  # 支持元数据
    AbbrExtension(),  # 支持缩写
    LegacyEmExtension(),  # 支持旧版强调
    MarkdownInHtmlExtension(),  # 支持 HTML 中的 Markdown
]

# 打开现有的 Whoosh 索引
try:
    ix = open_dir(index_dir)
except Exception as e:
    logging.debug(f"Error opening index: {e}")
    exit(1)


# 注册一个函数，在应用退出时关闭索引
def close_index():
    ix.close()


atexit.register(close_index)


def search_index(query_str):
    results = []
    try:
        if query_str.startswith("dir:"):
            # 提取目录路径
            dir_name = query_str[4:].strip()
            dir_path = os.path.join(BASE_DIR, dir_name)
            if os.path.isdir(dir_path):
                for root, _, files in os.walk(dir_path):
                    for file_name in files:
                        file_path = os.path.join(root, file_name)
                        folder_path = os.path.dirname(file_path)
                        results.append(
                            {
                                "file_name": file_name,
                                "file_path": file_path,
                                "folder_path": folder_path,
                                "score": 1,  # 假设所有文件的得分相同
                            }
                        )
            else:
                print(f"Directory not found: {dir_path}")
        else:
            with ix.searcher() as searcher:
                parser = QueryParser("file_content", ix.schema)
                query = parser.parse(query_str)
                for hit in searcher.search(query):
                    file_path = hit["file_path"]
                    file_name = hit["file_name"]
                    folder_path = os.path.dirname(file_path)
                    folder_name = os.path.basename(folder_path)
                    score = hit.score  # Assuming Whoosh provides a score
                    results.append(
                        {
                            "file_name": file_name,
                            "file_path": file_path,
                            "folder_path": folder_path,
                            "folder_name": folder_name,
                            "score": score,
                        }
                    )
    except Exception as e:
        print(f"Error during search: {e}")
    return results


app = Flask(__name__)


# 定义默认界面
@app.route("/")
def index():
    return render_template("index.html")


# 定义搜索路由
@app.route("/search", methods=["GET"])
def search():
    query_str = request.args.get("q")
    if not query_str:
        return redirect(url_for("index"))  # 重定向到根目录

    try:
        results = search_index(query_str)
        # return jsonify(results)
        return render_template("results.html", query=query_str, results=results)
    except Exception as e:
        return redirect(url_for("index")), 500


# 定义iframe默认页面路由
@app.route("/iframe_default")
def iframe_default():
    return render_template("iframe_default.html")


@app.route("/render_file")
def render_file():
    file_path = request.args.get("path")
    query_str = request.args.get("query")  # 获取搜索关键词
    if not file_path or not os.path.exists(file_path):
        return render_template("iframe_default.html"), 404

    try:
        file_extension = os.path.splitext(file_path)[1]
        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read()

        if file_extension == ".md":
            # 渲染 Markdown 文件
            rendered_content = markdown.markdown(
                content, extensions=markdown_extensions
            )

            # 高亮搜索关键词
            if query_str:
                soup = BeautifulSoup(rendered_content, "html.parser")
                pattern = re.compile(re.escape(query_str), re.IGNORECASE)

                # 遍历所有文本节点并高亮关键词
                for text_node in soup.find_all(text=True):
                    if text_node.parent.name not in [
                        "script",
                        "style",
                        "code",
                    ]:  # 避免处理脚本和样式
                        highlighted_text = pattern.sub(
                            f'<span style="background-color: yellow;">{query_str}</span>',
                            text_node,
                        )
                        text_node.replace_with(
                            BeautifulSoup(highlighted_text, "html.parser")
                        )

                rendered_content = str(soup)

            # 使用 render_template_string 函数来渲染包含 Jinja2 语法的字符串
            custom_css_link = "<link rel=\"stylesheet\" href=\"{{ url_for('static', filename='markdown_styles.css') }}\">"
            rendered_content = render_template_string(
                f"{custom_css_link}<div>{rendered_content}</div>", url_for=url_for
            )

        elif file_extension == ".html":
            # 直接返回 HTML 文件内容
            rendered_content = content

        elif file_extension == ".docx":
            # 渲染 DOCX 文件
            doc = docx.Document(file_path)
            rendered_content = "<html><body><h1>{}</h1>".format(
                doc.core_properties.title
            )
            for para in doc.paragraphs:
                rendered_content += "<p>{}</p>".format(para.text)
            rendered_content += "</body></html>"

        elif file_extension == ".pptx":
            # 渲染 PPT 文件
            prs = pptx.Presentation(file_path)
            rendered_content = "<html><body><h1>{}</h1>".format(
                prs.core_properties.title
            )
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        rendered_content += "<p>{}</p>".format(shape.text)
            rendered_content += "</body></html>"

        else:
            # 其他文件类型，直接返回文本内容
            rendered_content = f"<html><body><pre>{content}</pre></body></html>"

        # 添加滑动事件监测和 postMessage 传递
        rendered_content += """
        <script>
            let touchStartX = 0;
            let touchEndX = 0;

            document.addEventListener('touchstart', function (event) {
                touchStartX = event.changedTouches[0].screenX;
            });

            document.addEventListener('touchend', function (event) {
                touchEndX = event.changedTouches[0].screenX;
                handleSwipe();
            });

            function handleSwipe() {
                const swipeDistance = touchEndX - touchStartX;
                if (swipeDistance > 50) {
                    // Swipe right, show the list
                    window.parent.postMessage('showList', '*');
                } else if (swipeDistance < -50) {
                    // Swipe left, hide the list
                    window.parent.postMessage('hideList', '*');
                }
            }
        </script>
        """

        return rendered_content
    except Exception as e:
        return render_template("index.html"), 404


@app.route("/delete_file", methods=["DELETE"])
def delete_file():
    file_path = request.args.get("path")
    if not file_path:
        return jsonify({"error": "缺少文件路径"}), 400

    # 检查文件路径是否以任意一个允许的目录开头
    if not file_path.startswith(DEL_BASE_DIR):
        return jsonify({"error": "文件超出允许的目录"}), 400

    try:
        os.remove(file_path)
        logging.info(f"文件 '{file_path}' 已成功删除。")
        return jsonify({"message": "文件已成功删除"}), 200
    except FileNotFoundError:
        logging.error(f"删除文件 '{file_path}' 时出错: 找不到文件")
        return jsonify({"error": "找不到文件"}), 404
    except Exception as e:
        logging.error(f"删除文件 '{file_path}' 时出错: {e}")
        return jsonify({"error": "删除文件失败"}), 500


if __name__ == "__main__":
    app.run(debug=(log_level == logging.DEBUG))
